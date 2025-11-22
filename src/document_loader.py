
from pathlib import Path
from typing import List
import base64

from langchain_community.document_loaders import PyPDFLoader, UnstructuredPowerPointLoader
from langchain_core.documents import Document
from langchain_openai import ChatOpenAI

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    print("PyMuPDF not available. Image extraction from PDFs will be disabled.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("python-pptx not available. Image extraction from PowerPoint will be disabled.")


import hashlib
import concurrent.futures
from PIL import Image
import io

class DocumentLoader:    
    def __init__(self, api_key: str, vision_model: str = "x-ai/grok-4.1-fast"):
        self.api_key = api_key
        self.vision_llm = ChatOpenAI(
            model=vision_model,
            openai_api_key=api_key,
            openai_api_base="https://openrouter.ai/api/v1",
            temperature=0.2,
            max_tokens=4000
        )
    
    def load_pdf(self, pdf_path: str) -> List[Document]:
        """
        Load a PDF file and extract text
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            List of Document objects
        """
        try:
            loader = PyPDFLoader(pdf_path)
            documents = loader.load()
            
            for doc in documents:
                doc.metadata['source_file'] = Path(pdf_path).name
                doc.metadata['file_type'] = 'pdf'
            
            print(f"Loaded {len(documents)} pages from {Path(pdf_path).name}")
            return documents
            
        except Exception as e:
            print(f"Error loading PDF: {e}")
            return []
    
    def load_powerpoint(self, pptx_path: str) -> List[Document]:
        """
        Load a PowerPoint file and extract text
        
        Args:
            pptx_path: Path to the PowerPoint file
            
        Returns:
            List of Document objects
        """
        try:
            loader = UnstructuredPowerPointLoader(pptx_path)
            documents = loader.load()
            
            for doc in documents:
                doc.metadata['source_file'] = Path(pptx_path).name
                doc.metadata['file_type'] = 'pptx'
            
            print(f"Loaded PowerPoint: {Path(pptx_path).name}")
            return documents
            
        except Exception as e:
            print(f"Error loading PowerPoint: {e}")
            return []
    
    def extract_images_from_pdf(self, pdf_path: str, output_dir: str = "./temp_images") -> List[str]:
        """
        Extract images from PDF with filtering for small images and duplicates
        """
        if not PYMUPDF_AVAILABLE:
            print("PyMuPDF not installed. Cannot extract images.")
            return []
        
        Path(output_dir).mkdir(exist_ok=True)
        image_paths = []
        seen_hashes = set()
        
        try:
            pdf_document = fitz.open(pdf_path)
            
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                images = page.get_images()
                
                for img_index, img in enumerate(images):
                    xref = img[0]
                    base_image = pdf_document.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Filter small images and duplicates
                    try:
                        pil_img = Image.open(io.BytesIO(image_bytes))
                        width, height = pil_img.size
                        
                        # Skip small images (icons, bullets, etc.)
                        if width < 200 or height < 200:
                            continue
                            
                        # Skip duplicates
                        img_hash = hashlib.md5(image_bytes).hexdigest()
                        if img_hash in seen_hashes:
                            continue
                        seen_hashes.add(img_hash)
                        
                        image_path = Path(output_dir) / f"page{page_num + 1}_img{img_index + 1}.png"
                        with open(image_path, "wb") as img_file:
                            img_file.write(image_bytes)
                        
                        image_paths.append(str(image_path))
                        
                    except Exception as e:
                        continue
            
            pdf_document.close()
            print(f"Extracted {len(image_paths)} unique images from PDF (filtered small/duplicates)")
            return image_paths
            
        except Exception as e:
            print(f"Error extracting images: {e}")
            return []
    
    def extract_images_from_pptx(self, pptx_path: str, output_dir: str = "./temp_images") -> List[str]:
        """
        Extract images from PowerPoint with filtering
        """
        if not PPTX_AVAILABLE:
            print("python-pptx not installed. Cannot extract images.")
            return []
        
        Path(output_dir).mkdir(exist_ok=True)
        image_paths = []
        seen_hashes = set()
        
        try:
            prs = Presentation(pptx_path)
            
            for slide_num, slide in enumerate(prs.slides):
                for shape_num, shape in enumerate(slide.shapes):
                    if hasattr(shape, "image"):
                        image = shape.image
                        image_bytes = image.blob
                        
                        try:
                            pil_img = Image.open(io.BytesIO(image_bytes))
                            width, height = pil_img.size
                            
                            # Skip small images
                            if width < 200 or height < 200:
                                continue
                                
                            # Skip duplicates
                            img_hash = hashlib.md5(image_bytes).hexdigest()
                            if img_hash in seen_hashes:
                                continue
                            seen_hashes.add(img_hash)
                            
                            image_path = Path(output_dir) / f"slide{slide_num + 1}_img{shape_num + 1}.png"
                            with open(image_path, "wb") as img_file:
                                img_file.write(image_bytes)
                            
                            image_paths.append(str(image_path))
                            
                        except Exception:
                            continue
            
            print(f"Extracted {len(image_paths)} unique images from PowerPoint")
            return image_paths
            
        except Exception as e:
            print(f"Error extracting images: {e}")
            return []
    
    def analyze_image_with_vision(self, image_path: str) -> str:
        """
        Analyze an image using vision model
        """
        try:
            # Read and encode image
            with open(image_path, "rb") as img_file:
                image_bytes = img_file.read()
            
            base64_image = base64.b64encode(image_bytes).decode('utf-8')
            
            # Create vision prompt
            messages = [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Extract ONLY educational content (text, data, formulas). DO NOT describe colors, visual style, or layout."
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ]
            
            response = self.vision_llm.invoke(messages)
            return response.content
            
        except Exception as e:
            print(f"Error analyzing image {image_path}: {e}")
            return ""
    
    def process_images_to_docs(self, image_paths: List[str]) -> List[Document]:
        """
        Process images and create Document objects using parallel execution
        """
        image_docs = []
        
        print(f"Analyzing {len(image_paths)} images in parallel...")
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
            # Create a map of future -> image_path
            future_to_path = {executor.submit(self.analyze_image_with_vision, path): path for path in image_paths}
            
            for i, future in enumerate(concurrent.futures.as_completed(future_to_path), 1):
                path = future_to_path[future]
                try:
                    analysis = future.result()
                    if analysis:
                        doc = Document(
                            page_content=analysis,
                            metadata={
                                'source': path,
                                'type': 'image_analysis'
                            }
                        )
                        image_docs.append(doc)
                    print(f"Processed image {i}/{len(image_paths)}")
                except Exception as e:
                    print(f"Failed to process image {path}: {e}")
        
        print(f"Created {len(image_docs)} documents from images")
        return image_docs
